from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PRIMARY_BG = RGBColor(8, 22, 45)
SECONDARY_BG = RGBColor(24, 35, 62)
ACCENT_BLUE = RGBColor(0, 164, 255)
ACCENT_PURPLE = RGBColor(153, 102, 255)
ACCENT_TEAL = RGBColor(0, 206, 173)
ACCENT_ORANGE = RGBColor(255, 170, 92)
TEXT_LIGHT = RGBColor(243, 247, 255)
TEXT_DARK = RGBColor(24, 28, 40)


def style_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide, prs: Presentation) -> None:
    diameter = Inches(0.9)
    left = prs.slide_width - diameter - Inches(0.3)
    top = Inches(0.3)
    logo = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, diameter, diameter)
    logo.fill.solid()
    logo.fill.fore_color.rgb = ACCENT_TEAL
    logo.line.color.rgb = ACCENT_BLUE

    tf = logo.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "CJ"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.CENTER


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    style_background(slide, PRIMARY_BG)
    add_logo(slide, prs)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.color.rgb = TEXT_LIGHT

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.3), Inches(1.2), Inches(0.15), Inches(5.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_BLUE
    accent.line.width = 0

    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = TEXT_LIGHT


def add_picture_slide(prs: Presentation, title: str, image_path: Path, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_background(slide, SECONDARY_BG)
    add_logo(slide, prs)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT

    pic_left = Inches(0.7)
    pic_top = Inches(1.6)
    max_width = Inches(8.5)
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=max_width)
    else:
        placeholder = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, pic_left, pic_top, max_width, Inches(4.8)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = ACCENT_PURPLE
        placeholder.line.color.rgb = TEXT_LIGHT
        tf = placeholder.text_frame
        tf.text = "示意图缺失，请补充图片"
        tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    caption_box = slide.shapes.add_textbox(pic_left, pic_top + Inches(4.9), width=max_width, height=Inches(0.8))
    caption_tf = caption_box.text_frame
    caption_tf.clear()
    p = caption_tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER


def add_flow_slide(prs: Presentation, title: str, steps: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_background(slide, SECONDARY_BG)
    add_logo(slide, prs)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT

    left = Inches(0.5)
    top = Inches(1.7)
    box_width = Inches(2.2)
    box_height = Inches(1.1)
    gap = Inches(0.45)

    for idx, text in enumerate(steps):
        box_left = left + idx * (box_width + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, box_left, top, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE if idx % 2 == 0 else ACCENT_PURPLE
        shape.line.color.rgb = TEXT_LIGHT

        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER

        if idx > 0:
            arrow_left = box_left - gap / 1.6
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, arrow_left, top + box_height / 4, gap / 1.6, box_height / 2
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_TEAL
            arrow.line.width = 0


def add_layered_slide(prs: Presentation, title: str, layers: list[tuple[str, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_background(slide, PRIMARY_BG)
    add_logo(slide, prs)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT

    base_left = Inches(0.9)
    base_top = Inches(1.5)
    width = Inches(9.2)
    height = Inches(1.1)
    spacing = Inches(0.25)
    palette = [ACCENT_BLUE, ACCENT_PURPLE, ACCENT_TEAL, ACCENT_ORANGE]

    for idx, (layer_title, desc) in enumerate(layers):
        top = base_top + idx * (height + spacing)
        rect = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, base_left, top, width, height
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = palette[idx % len(palette)]
        rect.line.color.rgb = TEXT_LIGHT

        tf = rect.text_frame
        tf.clear()
        title_para = tf.paragraphs[0]
        title_para.text = layer_title
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.color.rgb = TEXT_DARK

        desc_para = tf.add_paragraph()
        desc_para.text = desc
        desc_para.font.size = Pt(16)
        desc_para.font.color.rgb = TEXT_DARK


def add_value_matrix_slide(prs: Presentation, title: str, matrix: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_background(slide, SECONDARY_BG)
    add_logo(slide, prs)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT

    rows = len(matrix)
    cols = len(matrix[0]) if rows else 0
    cell_width = Inches(3.1)
    cell_height = Inches(1.5)
    start_left = Inches(0.7)
    start_top = Inches(1.9)

    for r in range(rows):
        for c in range(cols):
            left = start_left + c * (cell_width + Inches(0.35))
            top = start_top + r * (cell_height + Inches(0.35))
            rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, cell_width, cell_height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = ACCENT_BLUE if (r + c) % 2 == 0 else ACCENT_PURPLE
            rect.line.color.rgb = TEXT_LIGHT
            tf = rect.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = matrix[r][c]
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER


def main() -> None:
    output_path = Path(__file__).with_name("MetaTransformation_plugin.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    style_background(title_slide, PRIMARY_BG)
    add_logo(title_slide, prs)
    title_slide.shapes.title.text = "仓颉 MetaTransformation 编译插件"
    subtitle = title_slide.placeholders[1]
    subtitle.text = "流程 | 架构 | 价值 | 场景"
    subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT

    add_bullet_slide(
        prs,
        "插件执行流程",
        [
            "输入：通过编译选项传入插件动态库路径",
            "加载：CompilerInstance 打开动态库并验证版本",
            "注册：插件向 Builder 注入 MetaTransform 实例",
            "执行：ToCHIR 对函数/包级 IR 执行 Run",
            "收尾：可选 IRCheck + Profile 审核变换质量",
        ],
    )

    add_layered_slide(
        prs,
        "架构分层与职责",
        [
            ("Loader 层", "动态库加载、版本校验、异常隔离"),
            ("Builder 层", "聚合插件注册回调，延迟构建变换"),
            ("Manager 层", "持有 MetaTransform 序列，统一调度"),
            ("执行管线", "在 CHIR 阶段遍历函数/包并执行"),
        ],
    )

    add_bullet_slide(
        prs,
        "核心数据结构",
        [
            "MetaTransformKind：描述函数级/包级插件类别",
            "MetaTransform<T>：模板抽象，约束 Run 接口",
            "MetaTransformPluginManager：遍历与调度容器",
            "MetaTransformPluginInfo：导出注册函数与版本",
            "CHIR_PLUGIN 宏：封装实例化与注册细节",
        ],
    )

    add_bullet_slide(
        prs,
        "实现方案与最佳实践",
        [
            "在插件内继承 MetaTransform<CHIR::Func/Package>",
            "构造函数注入 CHIRBuilder 等依赖并缓存",
            "registerTo 中调用 RegisterCHIRPluginCallback",
            "需要顺序依赖时按注册顺序组织回调",
            "启用 IRCheck/Profiling 监控插件质量",
        ],
    )

    add_bullet_slide(
        prs,
        "代码优势",
        [
            "强类型模板自动设置插件类别，避免误用",
            "Builder/Manager 解耦生命周期，扩展简单",
            "异常捕获与诊断隔离第三方故障",
            "可插拔设计保持主线编译器轻量",
        ],
    )

    add_bullet_slide(
        prs,
        "业务价值",
        [
            "面向行业的 IR 级优化与特性强化",
            "编译期静态审计、合规检测与安全插桩",
            "跨团队扩展与合作伙伴共创",
            "支持实验性功能快速验证与回滚",
        ],
    )

    add_bullet_slide(
        prs,
        "商业落地示例",
        [
            "金融风控 DSL：编译期降阶 + 合规审计",
            "车载安全：自动注入遥测与安全断言",
            "SaaS 多租户：租户级定制优化策略",
            "AIOps：编译期植入可观测性钩子",
        ],
    )

    figures_dir = Path(__file__).resolve().parents[1] / "figures"
    add_picture_slide(
        prs,
        "编译器系统架构",
        figures_dir / "Compiler_Architecture_Diagram_zh.png",
        "前端、LLVM 后端及 OS 支撑关系，标注插件嵌入点。",
    )
    add_picture_slide(
        prs,
        "单链路视角",
        figures_dir / "Compiler_Architecture_Diagram.png",
        "英文版架构示意，适合对外沟通或联合评审。",
    )

    add_flow_slide(
        prs,
        "MetaTransformation 插件管线",
        ["配置插件路径", "动态库加载 + 校验", "Builder 注册回调", "CHIR 执行变换", "IRCheck 验证"],
    )

    add_value_matrix_slide(
        prs,
        "价值映射矩阵",
        [
            ["性能增强\n(HPC/AI)", "安全合规\n(金融/车载)"],
            ["生态共建\n(伙伴插件)", "可观测 & AIOps\n(调试/监控)"],
        ],
    )

    prs.save(output_path)
    print(f"PPT saved to {output_path}")


if __name__ == "__main__":
    main()

